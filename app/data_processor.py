"""
Procesador de Datos Optimizado
===============================
- Limpieza y compresión de datos para maximizar densidad de información
- Procesamiento paralelizado para i7-9700K (8 núcleos)
- Recolección de basura explícita para liberar RAM inmediatamente
- Soporte completo: PDF, DOCX, XLSX, PPTX, imágenes (OCR), ZIP y más
"""
import gc
import json
import csv
import io
import os
import uuid
import zipfile
import tarfile
import tempfile
from pathlib import Path
from typing import Optional, Tuple
import pandas as pd
import pyarrow as pa
import pyarrow.parquet as pq

from app.config import (
    UPLOAD_CONFIG,
    TOKEN_ESTIMATION,
    HARDWARE_PROFILE,
    GC_CONFIG,
)


def estimate_tokens(text: str) -> int:
    """Estima tokens en tiempo real basado en chars_per_token.
    
    Para español con acentos y caracteres UTF-8 usamos 3.5 chars/token.
    """
    if not text:
        return 0
    return int(len(text) / TOKEN_ESTIMATION["chars_per_token"])


def optimize_dataset_for_llm(df: pd.DataFrame) -> str:
    """Convierte un DataFrame en texto plano compacto optimizado para LLM.
    
    Estrategias de compresión:
    1. Elimina NaN/None → ""
    2. Convierte a JSON compacto por fila (sin espacios redundantes)
    3. Cabecera descriptiva para contexto
    """
    # Limpieza rápida
    df = df.fillna("")
    
    # Convertir columnas a string optimizado
    for col in df.select_dtypes(include=["object"]).columns:
        df[col] = df[col].str.strip()
    
    # Generar representación compacta
    # Formato: cada fila como JSON compacto en una línea
    lines = []
    
    # Cabecera con metadatos
    lines.append(f"# DATASET: {len(df)} rows x {len(df.columns)} columns")
    lines.append(f"# COLUMNS: {', '.join(df.columns.tolist())}")
    lines.append("")
    
    # Datos compactos
    records = df.to_dict(orient="records")
    for record in records:
        # JSON compacto sin espacios
        line = json.dumps(record, ensure_ascii=False, separators=(",", ":"))
        lines.append(line)
    
    return "\n".join(lines)


def process_csv(filepath: str, optimize: bool = True) -> Tuple[str, int, dict]:
    """Procesa un archivo CSV y devuelve texto optimizado + metadatos.
    
    Args:
        filepath: Ruta al archivo CSV
        optimize: Si aplicar optimizaciones de compresión
        
    Returns:
        Tuple[str, int, dict]: (texto_optimizado, tokens_estimados, metadatos)
    """
    # Leer con detección automática de encoding
    try:
        df = pd.read_csv(filepath, encoding="utf-8")
    except UnicodeDecodeError:
        df = pd.read_csv(filepath, encoding="latin1")
    except Exception:
        # Fallback: intentar con pandas inferencia
        df = pd.read_csv(filepath)
    
    metadata = {
        "rows": len(df),
        "columns": len(df.columns),
        "column_names": df.columns.tolist(),
        "dtypes": {str(k): str(v) for k, v in df.dtypes.items()},
        "null_counts": df.isnull().sum().to_dict(),
        "memory_mb": round(df.memory_usage(deep=True).sum() / (1024 * 1024), 2),
    }
    
    if optimize:
        text = optimize_dataset_for_llm(df)
    else:
        text = df.to_string(index=False)
    
    tokens = estimate_tokens(text)

    del df
    return text, tokens, metadata


def process_json(filepath: str, optimize: bool = True) -> Tuple[str, int, dict]:
    """Procesa un archivo JSON/JSONL y devuelve texto optimizado + metadatos."""
    # Intentar como NDJSON primero
    try:
        df = pd.read_json(filepath, lines=True)
    except ValueError:
        # JSON estándar (array de objetos)
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            raw = f.read()
        data = json.loads(raw)
        if isinstance(data, list):
            df = pd.DataFrame(data)
        elif isinstance(data, dict):
            # Convertir dict a DataFrame si es posible
            df = pd.json_normalize(data)
        else:
            df = pd.DataFrame([data])
    
    metadata = {
        "rows": len(df),
        "columns": len(df.columns),
        "column_names": df.columns.tolist(),
        "memory_mb": round(df.memory_usage(deep=True).sum() / (1024 * 1024), 2),
    }
    
    if optimize:
        text = optimize_dataset_for_llm(df)
    else:
        text = df.to_string(index=False)
    
    tokens = estimate_tokens(text)

    del df
    return text, tokens, metadata


def process_text(filepath: str) -> Tuple[str, int, dict]:
    """Procesa un archivo de texto plano (.txt, .log)."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        raw_text = f.read()
    
    # Limpieza: eliminar líneas vacías múltiples, espacios redundantes
    lines = raw_text.split("\n")
    cleaned_lines = []
    prev_empty = False
    for line in lines:
        stripped = line.strip()
        if stripped == "":
            if not prev_empty:
                cleaned_lines.append("")
                prev_empty = True
        else:
            cleaned_lines.append(stripped)
            prev_empty = False
    
    text = "\n".join(cleaned_lines)
    tokens = estimate_tokens(text)
    
    metadata = {
        "original_chars": len(raw_text),
        "compressed_chars": len(text),
        "reduction_pct": round((1 - len(text) / max(len(raw_text), 1)) * 100, 1),
        "lines": len(cleaned_lines),
        "rows": len(cleaned_lines),
        "columns": 0,
    }
    
    del raw_text, lines, cleaned_lines
    return text, tokens, metadata


def truncate_to_token_limit(text: str, max_tokens: int = None) -> Tuple[str, int]:
    """Trunca el texto al límite de tokens permitido.
    
    Usa truncamiento inteligente: preserva cabeceras/metadatos primero.
    """
    if max_tokens is None:
        max_tokens = TOKEN_ESTIMATION["max_tokens"]
    
    current_tokens = estimate_tokens(text)
    if current_tokens <= max_tokens:
        return text, current_tokens
    
    # Preservar cabeceras (líneas que empiezan con #)
    lines = text.split("\n")
    header_lines = [l for l in lines if l.startswith("#")]
    data_lines = [l for l in lines if not l.startswith("#")]
    
    # Calcular cuánto espacio queda para datos después de cabeceras
    header_text = "\n".join(header_lines)
    header_tokens = estimate_tokens(header_text)
    available_tokens = max_tokens - header_tokens - 10  # 10 tokens de margen
    
    if available_tokens <= 0:
        # Solo devolver cabeceras
        return header_text, header_tokens
    
    # Truncar datos línea por línea
    truncated_data = []
    data_tokens = 0
    for line in data_lines:
        line_tokens = estimate_tokens(line) + 1  # +1 por newline
        if data_tokens + line_tokens <= available_tokens:
            truncated_data.append(line)
            data_tokens += line_tokens
        else:
            break
    
    # Reconstruir texto
    result = header_text + "\n" + "\n".join(truncated_data)
    final_tokens = estimate_tokens(result)
    
    del lines, header_lines, data_lines
    return result, final_tokens


def _try_import(module: str):
    """Importa un módulo de forma segura; retorna None si no está instalado."""
    import importlib
    try:
        return importlib.import_module(module)
    except ImportError:
        return None


def process_pdf(filepath: str) -> Tuple[str, int, dict]:
    """Extrae texto de un PDF usando pypdf o pdfplumber como fallback."""
    text_parts = []
    num_pages = 0

    # Intentar con pypdf primero
    pypdf = _try_import("pypdf")
    if pypdf:
        try:
            reader = pypdf.PdfReader(filepath)
            num_pages = len(reader.pages)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"[Página {i+1}]\n{page_text}")
        except Exception as e:
            text_parts = []

    # Fallback: pdfplumber
    if not text_parts:
        pdfplumber = _try_import("pdfplumber")
        if pdfplumber:
            try:
                with pdfplumber.open(filepath) as pdf:
                    num_pages = len(pdf.pages)
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            text_parts.append(f"[Página {i+1}]\n{page_text}")
            except Exception:
                pass

    if not text_parts:
        text = f"[PDF: {Path(filepath).name} — {num_pages} páginas. No se pudo extraer texto (posiblemente escaneado o protegido).]"
    else:
        text = f"[PDF: {num_pages} páginas]\n\n" + "\n\n".join(text_parts)

    tokens = estimate_tokens(text)
    metadata = {"pages": num_pages, "extracted_pages": len(text_parts)}
    return text, tokens, metadata


def process_docx(filepath: str) -> Tuple[str, int, dict]:
    """Extrae texto de un archivo Word .docx."""
    docx = _try_import("docx")
    if docx is None:
        # Fallback: python-docx se importa como 'docx'
        try:
            import docx as docx_module
        except ImportError:
            text = f"[DOCX: {Path(filepath).name} — módulo python-docx no instalado]"
            return text, estimate_tokens(text), {"error": "missing python-docx"}
        docx = docx_module
    try:
        import docx as docx_module
        doc = docx_module.Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Extraer tablas también
        table_texts = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            table_texts.append("\n".join(rows))
        full_text = "\n".join(paragraphs)
        if table_texts:
            full_text += "\n\n[TABLAS]\n" + "\n\n".join(table_texts)
        metadata = {"paragraphs": len(paragraphs), "tables": len(doc.tables)}
        return full_text, estimate_tokens(full_text), metadata
    except Exception as e:
        text = f"[DOCX: Error al procesar {Path(filepath).name}: {e}]"
        return text, estimate_tokens(text), {"error": str(e)}


def process_xlsx(filepath: str) -> Tuple[str, int, dict]:
    """Extrae datos de un Excel .xlsx/.xls/.ods."""
    ext = Path(filepath).suffix.lower()
    try:
        if ext == ".xls":
            dfs = pd.read_excel(filepath, engine="xlrd", sheet_name=None)
        elif ext == ".ods":
            dfs = pd.read_excel(filepath, engine="odf", sheet_name=None)
        else:
            dfs = pd.read_excel(filepath, sheet_name=None)

        parts = []
        sheet_summaries = []
        for sheet_name, df in dfs.items():
            df = df.fillna("")
            parts.append(f"[Hoja: {sheet_name} — {len(df)} filas × {len(df.columns)} columnas]\n{optimize_dataset_for_llm(df)}")
            sheet_summaries.append({"name": sheet_name, "rows": len(df), "cols": len(df.columns)})
        text = "\n\n".join(parts)
        metadata = {"sheets": sheet_summaries, "total_sheets": len(dfs)}
        return text, estimate_tokens(text), metadata
    except Exception as e:
        text = f"[Excel: Error al procesar {Path(filepath).name}: {e}]"
        return text, estimate_tokens(text), {"error": str(e)}


def process_pptx(filepath: str) -> Tuple[str, int, dict]:
    """Extrae texto de un PowerPoint .pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                slides_text.append(f"[Diapositiva {i}]\n" + "\n".join(slide_parts))
        text = f"[PPTX: {len(prs.slides)} diapositivas]\n\n" + "\n\n".join(slides_text)
        metadata = {"slides": len(prs.slides), "slides_with_text": len(slides_text)}
        return text, estimate_tokens(text), metadata
    except ImportError:
        text = f"[PPTX: {Path(filepath).name} — módulo python-pptx no instalado]"
        return text, estimate_tokens(text), {"error": "missing python-pptx"}
    except Exception as e:
        text = f"[PPTX: Error al procesar {Path(filepath).name}: {e}]"
        return text, estimate_tokens(text), {"error": str(e)}


def process_image(filepath: str) -> Tuple[str, int, dict]:
    """Análisis profundo de imagen: metadatos, colores dominantes, OCR y base64 para visión IA."""
    metadata = {"type": "image"}
    parts = [f"[IMAGEN: {Path(filepath).name}]"]

    PIL = _try_import("PIL.Image")
    if PIL:
        try:
            from PIL import Image, ImageStat
            img = Image.open(filepath)
            w, h = img.size
            fmt = img.format or Path(filepath).suffix.upper().lstrip(".")
            mode = img.mode
            parts.append(f"Dimensiones: {w}×{h}px | Formato: {fmt} | Modo: {mode}")
            parts.append(f"Píxeles totales: {w*h:,} | Aspecto: {w/h:.2f}:1")
            metadata.update({"width": w, "height": h, "mode": mode, "format": fmt})

            # Estadísticas de color por canal
            img_rgb = img.convert("RGB")
            stat = ImageStat.Stat(img_rgb)
            r_mean, g_mean, b_mean = stat.mean[:3]
            brightness = r_mean * 0.299 + g_mean * 0.587 + b_mean * 0.114
            brightness_label = "oscura" if brightness < 85 else "media" if brightness < 170 else "clara"
            parts.append(f"Brillo: {brightness:.0f}/255 ({brightness_label}) | RGB medio: ({r_mean:.0f}, {g_mean:.0f}, {b_mean:.0f})")

            # Colores dominantes (cuantización adaptativa a 8 colores)
            try:
                img_small = img_rgb.resize((150, 150))
                result = img_small.convert("P", palette=Image.ADAPTIVE, colors=8)
                palette = result.getpalette()
                dominant = [f"#{palette[i*3]:02x}{palette[i*3+1]:02x}{palette[i*3+2]:02x}" for i in range(8)]
                parts.append(f"Colores dominantes: {', '.join(dominant)}")
                metadata["dominant_colors"] = dominant
            except Exception:
                pass

            # EXIF si está disponible
            try:
                from PIL.ExifTags import TAGS
                exif_raw = getattr(img, "_getexif", lambda: None)()
                if exif_raw:
                    useful = {TAGS.get(k, k): str(v) for k, v in exif_raw.items()
                              if TAGS.get(k, k) in ("Make", "Model", "DateTime", "Software", "ImageDescription")}
                    if useful:
                        parts.append("EXIF: " + " | ".join(f"{k}={v}" for k, v in useful.items()))
            except Exception:
                pass

        except Exception as e:
            parts.append(f"[Error PIL: {e}]")
    else:
        parts.append("[PIL/Pillow no disponible — instala Pillow para análisis de imágenes]")

    # OCR con pytesseract
    pytesseract = _try_import("pytesseract")
    if pytesseract:
        try:
            from PIL import Image
            import pytesseract as tess
            img_ocr = Image.open(filepath)
            ocr_text = tess.image_to_string(img_ocr, lang="spa+eng").strip()
            if ocr_text:
                parts.append(f"\n[TEXTO DETECTADO POR OCR]\n{ocr_text}")
                metadata["ocr_chars"] = len(ocr_text)
                metadata["has_text"] = True
        except Exception:
            pass

    # Codificar imagen en base64 para análisis visual con modelo de visión
    try:
        import base64
        with open(filepath, "rb") as f:
            raw_bytes = f.read()
        img_b64 = base64.b64encode(raw_bytes).decode("utf-8")
        ext_lower = Path(filepath).suffix.lower()
        mime_map = {
            ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp",
            ".tiff": "image/tiff", ".tif": "image/tiff",
        }
        metadata["_vision_b64"] = img_b64
        metadata["_vision_mime"] = mime_map.get(ext_lower, "image/png")
        metadata["has_vision_data"] = True
        parts.append("\n[Imagen lista para análisis visual por IA — se enviará al modelo]")
    except Exception:
        pass

    text = "\n".join(parts)
    return text, estimate_tokens(text), metadata


def process_svg(filepath: str) -> Tuple[str, int, dict]:
    """Parsea un archivo SVG extrayendo estructura visual, colores y texto."""
    import xml.etree.ElementTree as ET

    parts = [f"[SVG: {Path(filepath).name}]"]
    metadata = {"type": "svg"}

    try:
        tree = ET.parse(filepath)
        root = tree.getroot()

        # Namespace
        SVG_NS = "http://www.w3.org/2000/svg"
        def tag_name(elem):
            t = elem.tag
            return t.split("}")[-1] if "}" in t else t

        # Atributos raíz
        vb = root.get("viewBox", "")
        w = root.get("width", "")
        h = root.get("height", "")
        if vb:
            parts.append(f"ViewBox: {vb}")
        if w or h:
            parts.append(f"Dimensiones: {w} × {h}")
        metadata.update({"viewBox": vb, "width": w, "height": h})

        # Inventario de elementos y colores
        elem_count: dict = {}
        colors: set = set()
        texts_found: list = []

        for elem in root.iter():
            tname = tag_name(elem)
            elem_count[tname] = elem_count.get(tname, 0) + 1

            # Colores de atributos y estilos inline
            for attr in ("fill", "stroke"):
                val = elem.get(attr, "")
                if val and val not in ("none", "transparent", "inherit", ""):
                    colors.add(val)
            style = elem.get("style", "")
            for prop in style.split(";"):
                if ":" in prop:
                    k, v = prop.split(":", 1)
                    if k.strip() in ("fill", "stroke") and v.strip() not in ("none", "transparent", ""):
                        colors.add(v.strip())

            # Texto
            if tname in ("text", "tspan", "textPath"):
                content = "".join(elem.itertext()).strip()
                if content:
                    texts_found.append(content)

        # Resumen elementos (sin svg/g/defs)
        visual = {k: v for k, v in elem_count.items()
                  if k not in ("svg", "g", "defs", "style", "clipPath", "mask", "pattern")}
        if visual:
            top = sorted(visual.items(), key=lambda x: -x[1])[:12]
            parts.append(f"Formas: {', '.join(f'{k}×{v}' for k, v in top)}")

        if colors:
            parts.append(f"Colores usados: {', '.join(sorted(colors)[:20])}")
            metadata["colors"] = sorted(colors)[:20]

        if texts_found:
            parts.append(f"Textos en el SVG:")
            for t in texts_found[:30]:
                parts.append(f"  · {t[:300]}")
            metadata["texts"] = texts_found[:30]

        # title / desc del SVG
        for tag in ("title", "desc"):
            el = root.find(f"{{{SVG_NS}}}{tag}")
            if el is None:
                el = root.find(tag)
            if el is not None and el.text:
                parts.append(f"{tag.capitalize()}: {el.text.strip()}")

        # Código fuente completo (los SVG suelen ser texto legible)
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            raw_svg = f.read()
        max_svg_chars = 40_000
        if len(raw_svg) <= max_svg_chars:
            parts.append(f"\n[CÓDIGO SVG]\n{raw_svg}")
        else:
            parts.append(f"\n[CÓDIGO SVG (primeros {max_svg_chars:,} chars de {len(raw_svg):,}])\n{raw_svg[:max_svg_chars]}")

        metadata.update({"elem_count": elem_count, "total_elements": sum(visual.values())})

    except ET.ParseError as e:
        parts.append(f"[Error al parsear SVG como XML: {e}]")
        # Fallback: leer como texto
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            parts.append(f.read()[:30_000])
    except Exception as e:
        parts.append(f"[Error SVG: {e}]")

    text = "\n".join(parts)
    return text, estimate_tokens(text), metadata


def process_zip(filepath: str) -> Tuple[str, int, dict]:
    """Extrae y analiza el contenido de un archivo ZIP.
    
    - Lista todos los archivos del ZIP con tamaños
    - Extrae y procesa los archivos de texto/código automáticamente
    - Procesa recursivamente ZIPs anidados
    - Omite binarios y archivos demasiado grandes
    """
    MAX_EXTRACT_SIZE = 50 * 1024 * 1024  # 50 MB por archivo
    MAX_TOTAL_TEXT = 200_000  # chars máximos de texto total
    SKIP_EXTENSIONS = {".pyc", ".pyo", ".class", ".o", ".so", ".dll", ".exe",
                       ".bin", ".dat", ".pack", ".idx"}

    parts = [f"[ARCHIVO ZIP: {Path(filepath).name}]"]
    file_inventory = []
    processed_texts = []
    errors = []

    try:
        with zipfile.ZipFile(filepath, "r") as zf:
            members = zf.infolist()
            parts.append(f"Total de entradas: {len(members)}")

            # Inventario
            for info in members:
                if not info.is_dir():
                    size_kb = info.file_size / 1024
                    file_inventory.append(f"  {info.filename} ({size_kb:.1f} KB)")

            parts.append("\n[INVENTARIO DE ARCHIVOS]\n" + "\n".join(file_inventory[:200]))
            if len(file_inventory) > 200:
                parts.append(f"  ... y {len(file_inventory) - 200} archivos más")

            # Extraer y procesar archivos de texto
            total_chars = 0
            with tempfile.TemporaryDirectory() as tmpdir:
                for info in members:
                    if info.is_dir():
                        continue
                    ext = Path(info.filename).suffix.lower()
                    if ext in SKIP_EXTENSIONS:
                        continue
                    if info.file_size > MAX_EXTRACT_SIZE:
                        processed_texts.append(f"\n[{info.filename}] — omitido (demasiado grande: {info.file_size/1024/1024:.1f} MB)")
                        continue
                    if total_chars >= MAX_TOTAL_TEXT:
                        processed_texts.append(f"\n[... límite de texto alcanzado, {len(members)} archivos en total]")
                        break

                    try:
                        data = zf.read(info.filename)
                        # Si es otro ZIP, procesar recursivamente
                        if ext == ".zip":
                            sub_path = os.path.join(tmpdir, Path(info.filename).name)
                            with open(sub_path, "wb") as f:
                                f.write(data)
                            sub_text, _, sub_meta = process_zip(sub_path)
                            processed_texts.append(f"\n[ZIP ANIDADO: {info.filename}]\n{sub_text[:5000]}")
                            total_chars += len(sub_text[:5000])
                            continue

                        # Intentar decodificar como texto
                        file_text = None
                        for enc in ("utf-8", "latin-1", "cp1252"):
                            try:
                                file_text = data.decode(enc)
                                break
                            except Exception:
                                continue

                        if file_text is None:
                            # Archivo binario — saltar
                            continue

                        # Truncar si es muy largo
                        if len(file_text) > 20000:
                            file_text = file_text[:20000] + f"\n... [truncado, {len(data)} bytes totales]"

                        processed_texts.append(f"\n{'='*50}\n[ARCHIVO: {info.filename}]\n{'='*50}\n{file_text}")
                        total_chars += len(file_text)

                    except Exception as e:
                        errors.append(f"{info.filename}: {e}")

        if processed_texts:
            parts.append("\n[CONTENIDO DE ARCHIVOS]")
            parts.extend(processed_texts)
        if errors:
            parts.append(f"\n[ERRORES AL PROCESAR: {', '.join(errors[:10])}]")

        metadata = {
            "total_files": len(file_inventory),
            "processed_files": len(processed_texts),
            "type": "zip",
        }

    except zipfile.BadZipFile:
        parts.append("[ERROR: Archivo ZIP dañado o inválido]")
        metadata = {"error": "bad_zip"}
    except Exception as e:
        parts.append(f"[ERROR al abrir ZIP: {e}]")
        metadata = {"error": str(e)}

    text = "\n".join(parts)
    return text, estimate_tokens(text), metadata


def process_tar(filepath: str) -> Tuple[str, int, dict]:
    """Extrae y analiza un archivo TAR (.tar, .tar.gz, .tgz, .tar.bz2, .tar.xz)."""
    MAX_EXTRACT_SIZE = 50 * 1024 * 1024
    MAX_TOTAL_TEXT = 200_000
    SKIP_EXTENSIONS = {".pyc", ".pyo", ".class", ".o", ".so", ".dll", ".exe",
                       ".bin", ".dat", ".pack", ".idx"}

    parts = [f"[ARCHIVO TAR: {Path(filepath).name}]"]
    file_inventory = []
    processed_texts = []

    try:
        with tarfile.open(filepath, "r:*") as tf:
            members = tf.getmembers()
            files = [m for m in members if m.isfile()]
            parts.append(f"Total de archivos: {len(files)}")

            for m in files:
                file_inventory.append(f"  {m.name} ({m.size/1024:.1f} KB)")
            parts.append("\n[INVENTARIO]\n" + "\n".join(file_inventory[:200]))

            total_chars = 0
            for m in files:
                if total_chars >= MAX_TOTAL_TEXT:
                    break
                ext = Path(m.name).suffix.lower()
                if ext in SKIP_EXTENSIONS or m.size > MAX_EXTRACT_SIZE:
                    continue
                try:
                    f = tf.extractfile(m)
                    if f is None:
                        continue
                    data = f.read()
                    file_text = None
                    for enc in ("utf-8", "latin-1"):
                        try:
                            file_text = data.decode(enc)
                            break
                        except Exception:
                            continue
                    if file_text is None:
                        continue
                    if len(file_text) > 15000:
                        file_text = file_text[:15000] + "\n... [truncado]"
                    processed_texts.append(f"\n{'='*50}\n[{m.name}]\n{'='*50}\n{file_text}")
                    total_chars += len(file_text)
                except Exception:
                    continue

        if processed_texts:
            parts.append("\n[CONTENIDO]")
            parts.extend(processed_texts)

        metadata = {"total_files": len(files), "type": "tar"}
    except Exception as e:
        parts.append(f"[ERROR al abrir TAR: {e}]")
        metadata = {"error": str(e)}

    text = "\n".join(parts)
    return text, estimate_tokens(text), metadata


def process_ipynb(filepath: str) -> Tuple[str, int, dict]:
    """Extrae celdas de código y markdown de un Jupyter Notebook."""
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            nb = json.load(f)

        cells = nb.get("cells", [])
        parts = [f"[Jupyter Notebook: {len(cells)} celdas]"]
        code_cells = 0
        md_cells = 0

        for i, cell in enumerate(cells, 1):
            ctype = cell.get("cell_type", "")
            source = "".join(cell.get("source", []))
            if not source.strip():
                continue
            if ctype == "code":
                code_cells += 1
                parts.append(f"\n[Celda {i} — Código]\n```python\n{source}\n```")
                # Incluir outputs de texto si los hay
                outputs = cell.get("outputs", [])
                for out in outputs:
                    if out.get("output_type") in ("stream", "execute_result", "display_data"):
                        out_text = "".join(out.get("text", out.get("data", {}).get("text/plain", [])))
                        if out_text.strip():
                            parts.append(f"[Output]: {out_text[:500]}")
            elif ctype == "markdown":
                md_cells += 1
                parts.append(f"\n[Celda {i} — Markdown]\n{source}")

        text = "\n".join(parts)
        metadata = {"total_cells": len(cells), "code_cells": code_cells, "md_cells": md_cells}
        return text, estimate_tokens(text), metadata
    except Exception as e:
        text = f"[ipynb: Error al procesar {Path(filepath).name}: {e}]"
        return text, estimate_tokens(text), {"error": str(e)}


def process_file(filepath: str) -> Tuple[str, int, dict]:
    """Dispatcher universal: procesa cualquier tipo de archivo.

    Soporta:
    - Datos: CSV, TSV, JSON, JSONL → procesamiento estructurado
    - Office: PDF, DOCX, DOC, XLSX, XLS, ODS, PPTX, PPT → extracción de contenido
    - Imágenes: PNG, JPG, WEBP, etc. → metadatos + OCR opcional
    - Archivos comprimidos: ZIP, TAR, GZ, TGZ → inventario + extracción de texto
    - Notebooks: IPYNB → celdas de código y markdown
    - Texto / Código: todo lo demás → lectura como texto plano
    """
    ext = Path(filepath).suffix.lower()
    name = Path(filepath).name

    # ── Datos estructurados ──
    if ext == ".csv":
        return process_csv(filepath)
    elif ext in (".json", ".jsonl"):
        return process_json(filepath)
    elif ext == ".tsv":
        df = pd.read_csv(filepath, sep="\t")
        metadata = {"rows": len(df), "columns": len(df.columns), "column_names": df.columns.tolist()}
        text = optimize_dataset_for_llm(df)
        tokens = estimate_tokens(text)
        del df
        return text, tokens, metadata

    # ── Documentos Office / PDF ──
    elif ext == ".pdf":
        return process_pdf(filepath)
    elif ext in (".docx", ".doc"):
        return process_docx(filepath)
    elif ext in (".xlsx", ".xls", ".xlsm", ".ods"):
        return process_xlsx(filepath)
    elif ext in (".pptx", ".ppt"):
        return process_pptx(filepath)

    # ── Imágenes raster ──
    elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"):
        return process_image(filepath)

    # ── Imágenes vectoriales ──
    elif ext == ".svg":
        return process_svg(filepath)

    # ── Archivos comprimidos ──
    elif ext == ".zip":
        return process_zip(filepath)
    elif ext in (".tar", ".gz", ".tgz", ".tar.gz", ".tar.bz2", ".tar.xz"):
        # .gz podría ser un tar.gz o un gz simple
        if ext == ".gz" and not name.endswith(".tar.gz"):
            # gz simple: descomprimir y tratar como texto
            import gzip
            try:
                with gzip.open(filepath, "rt", encoding="utf-8", errors="replace") as f:
                    content = f.read(100_000)
                metadata = {"type": "gzip", "size": len(content)}
                text = f"[GZ: {name}]\n{content}"
                return text, estimate_tokens(text), metadata
            except Exception as e:
                text = f"[GZ: Error al descomprimir {name}: {e}]"
                return text, estimate_tokens(text), {"error": str(e)}
        return process_tar(filepath)

    # ── Jupyter Notebook ──
    elif ext == ".ipynb":
        return process_ipynb(filepath)

    # ── Texto plano y código ──
    elif ext in (".txt", ".log", ".md", ".rst", ".rtf", ".odt", ".epub"):
        return process_text(filepath)

    # ── Cualquier otro (código, config, etc.) → texto plano ──
    else:
        return process_text(filepath)


def save_upload(file_data: bytes, original_filename: str) -> str:
    """Guarda un archivo subido en el directorio temporal."""
    os.makedirs(UPLOAD_CONFIG["temp_dir"], exist_ok=True)
    
    # Generar nombre único para evitar colisiones
    ext = Path(original_filename).suffix.lower()
    unique_name = f"{uuid.uuid4().hex}{ext}"
    filepath = os.path.join(UPLOAD_CONFIG["temp_dir"], unique_name)
    
    with open(filepath, "wb") as f:
        f.write(file_data)
    
    return filepath


def cleanup_file(filepath: str):
    """Elimina un archivo temporal y fuerza un único ciclo de GC por subida."""
    try:
        if os.path.exists(filepath):
            os.remove(filepath)
    except Exception:
        pass
    gc.collect()  # Un solo gc.collect() por ciclo completo de subida